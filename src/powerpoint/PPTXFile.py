import os
import re
import shutil
import tempfile
from functools import cached_property

import win32com.client
from gtts import gTTS
from moviepy.editor import (AudioFileClip, CompositeAudioClip, ImageClip,
                            VideoFileClip, afx, concatenate_videoclips)
from pydub import AudioSegment
from utils import Hash, Log

from pptx import Presentation as PPTXPresentation

log = Log('PPTXFile')


class PPTXFile:
    DELIM_NOTES = '\n'

    def __init__(self, file_path):
        self.file_path = file_path

    @cached_property
    def presentation(self) -> PPTXPresentation:
        return PPTXPresentation(self.file_path)

    @cached_property
    def dir_path(self):
        file_name_only = os.path.basename(self.file_path).split('.')[0]
        dir_path = os.path.join(
            tempfile.gettempdir(), f'pptx-{file_name_only}'
        )
        os.makedirs(dir_path, exist_ok=True)
        return dir_path

    @cached_property
    def image_path_list(self) -> list[str]:
        app = win32com.client.Dispatch("Powerpoint.Application")
        presentation = app.Presentations.Open(self.file_path)
        image_path_list = []
        for i, slide in enumerate(presentation.Slides):
            image_path = os.path.join(self.dir_path, f'{i:03d}.png')
            slide.Export(image_path, "PNG")
            image_path_list.append(image_path)
            log.debug(f'Wrote {image_path}')
        app.Quit()
        return image_path_list

    @staticmethod
    def clean_content(x: str) -> str:
        x = x.replace('AI', 'A.I.')
        x = x.replace('...', ' ')
        x = x.replace('..', ' ')
        x = x.replace('…', ' ')
        x = re.sub(r' +', ' ', x)
        return x

    @cached_property
    def notes_list(self) -> list[str]:
        notes_list = []
        for slide in self.presentation.slides:
            notes_content = PPTXFile.clean_content(
                slide.notes_slide.notes_text_frame.text
            )
            notes = notes_content.split(PPTXFile.DELIM_NOTES)
            # filter out links
            notes = [note for note in notes if 'http' not in note]
            notes_list.append(notes)
        return notes_list

    @staticmethod
    def get_audio_clip(
        path_base, notes: list[str], is_first: bool, is_last: bool
    ):
        content = '\n'.join(notes) + '\n\n'
        audio_path = path_base + '.mp3'
        if not os.path.exists(audio_path):
            tts = gTTS(content, lang='en', slow=False)
            tts.save(audio_path)
            log.debug(f'Wrote {audio_path}')

        start_duration = 2_000 if is_first else 500
        audio = AudioSegment.silent(
            duration=start_duration
        ) + AudioSegment.from_file(audio_path).speedup(playback_speed=1.2)

        end_duration = 10_000 if is_last else 500
        audio += AudioSegment.silent(duration=end_duration)
        audio.export(audio_path, format='mp3')

        audio_clip = AudioFileClip(audio_path)
        return audio_clip

    @staticmethod
    def get_video_clip(path_base, image_path: str, audio_clip: AudioFileClip):
        video_path = path_base + '.mp4'
        if not os.path.exists(video_path):
            clip = (
                ImageClip(image_path)
                .set_duration(audio_clip.duration)
                .set_audio(audio_clip)
            )
            clip.write_videofile(video_path, fps=24)
            log.debug(f'Wrote {video_path}')

        video_clip = VideoFileClip(video_path)
        return video_clip

    def write_video(self):
        video_clips = []

        for i, (notes, image_path) in enumerate(
            zip(self.notes_list, self.image_path_list)
        ):
            log.debug(f'{i=}, {notes=}')
            content = ' '.join(notes)
            h = Hash.md5(content)[:6]
            path_base = os.path.join(self.dir_path, h)

            is_first = i == 0
            is_last = i == len(self.notes_list) - 1
            audio_clip = PPTXFile.get_audio_clip(
                path_base, notes, is_first, is_last
            )
            video_clip = PPTXFile.get_video_clip(
                path_base, image_path, audio_clip
            )
            video_clips.append(video_clip)

        combined_video_clip = concatenate_videoclips(
            video_clips, method="compose"
        )

        audio_clip = afx.audio_loop(
            AudioFileClip(os.path.join('media', 'thelounge.mp3')).volumex(
                0.5
            ),
            duration=combined_video_clip.duration,
        )

        combined_audio_clip = CompositeAudioClip(
            [combined_video_clip.audio, audio_clip]
        )
        combined_video_clip.audio = combined_audio_clip

        combined_video_path = os.path.join(self.dir_path, 'video.mp4')
        combined_video_clip.write_videofile(
            combined_video_path,
            codec='mpeg4',
            fps=24,
            audio_codec='libmp3lame',
        )
        log.info(f'Wrote {combined_video_path}')

        copy_video_path = self.file_path.replace('.pptx', '-video.mp4')
        shutil.copy(combined_video_path, copy_video_path)
        log.info(f'Copied to {copy_video_path}')

        return combined_video_path
