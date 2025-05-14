from dataclasses import dataclass

from PIL import Image
from utils import File, Log

from powerpoint.script.PPTXScriptSlide import PPTXScriptSlide
from pptx import Presentation

log = Log("PPTXScript")


@dataclass
class PPTXScript:
    slides: list[PPTXScriptSlide]

    def write(self, pptx_path: str):
        prs = Presentation()

        for slide in self.slides:
            slide_layout = prs.slide_layouts[6]  # blank slide
            prs_slide = prs.slides.add_slide(slide_layout)

            assert len(slide.images) == 1

            with Image.open(slide.images[0]) as img:
                image_width, image_height = img.size

            slide_width, slide_height = prs.slide_width, prs.slide_height

            r = max(
                image_width / slide_width,
                image_height / slide_height,
            )
            image_display_width, image_display_height = (
                image_width / r,
                image_height / r,
            )

            padding_x = (slide_width - image_display_width) / 2
            padding_y = (slide_height - image_display_height) / 2
            padding = 500_000
            prs_slide.shapes.add_picture(
                slide.images[0],
                padding_x + padding,
                padding_y + padding,
                width=image_display_width - padding * 2,
                height=image_display_height - padding * 2,
            )

            if slide.notes:
                notes_slide = prs_slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = "\n\n".join(slide.notes)

        prs.save(pptx_path)
        log.info(f"Wrote {pptx_path}")

    def from_simple_config(config_lines):
        slides = []
        for line in config_lines:
            image_path, notes = line
            slides.append(
                PPTXScriptSlide(
                    text="",
                    images=[image_path],
                    notes=[notes],
                )
            )
        return PPTXScript(slides)

    def from_md(md_path):
        lines = File(md_path).read_lines()
        i = 0
        slides = []
        while i < len(lines):
            line = lines[i].strip()
            if not line:
                i += 1
                continue
            if line.startswith("#"):
                i += 1
                continue

            if line.startswith("!["):
                image_path = line.split("(")[1].split(")")[0]

                notes_lines = []
                i += 1

                while i < len(lines):
                    line = lines[i].strip()
                    if line.startswith("!["):
                        i -= 1
                        break
                    notes_lines.append(line)
                    i += 1

                notes = "\n".join(notes_lines)
                slides.append(
                    PPTXScriptSlide(
                        text="",
                        images=[image_path],
                        notes=[notes],
                    )
                )
                i += 1

        return PPTXScript(slides)
